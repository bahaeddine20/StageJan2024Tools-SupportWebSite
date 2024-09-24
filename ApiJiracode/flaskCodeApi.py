import base64
import datetime
import io
import json
import tempfile
import zipfile
import os
import matplotlib.pyplot as plt
import mysql.connector
import numpy as np
import pandas as pd
from flask import Flask, request
from flask_cors import CORS
from jira import JIRA
from mysql.connector import Error
from flask import request, send_file, abort
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
from flask import request
from openpyxl import load_workbook
from openpyxl.styles import PatternFill, Font, Border, Alignment
import xlwings as xw
from numpy import save
from openpyxl import Workbook
from openpyxl.chart import PieChart3D, Reference
from openpyxl.chart.label import DataLabelList
from openpyxl.chart.marker import DataPoint


pat = ""  #  When running locally
#pat = os.getenv('PAT')   #  When using Docker



#serverdb='mysqldb'  #When using Docker
#test

serverdb='localhost'  # When running locally





# Specify Jira connection parameters
jiraOptions = {'server': "https://jira.dt.renault.com"}
jira = JIRA(options=jiraOptions, token_auth=pat)
def get_board_id_by_project(project_key):
    """Retrieve the board ID associated with a project key."""
    boards = jira.boards()
    for board in boards:
        if board.location.projectKey == project_key:
            return board.id
    return None

def get_sprints(board_id):
    """Retrieve all sprints for the specified board."""
    sprints = jira.sprints(board_id)
    return sprints

def search_issues_by_sprint(sprint_id):
    """Search for all issues in a specified sprint."""
    jql_str = f'sprint = {sprint_id}'
    issues = jira.search_issues(jql_str, expand='changelog')
    return issues


def is_weekend(dt):
    """Check if the given datetime is on a weekend (Saturday or Sunday)."""
    return dt.weekday() >= 5  # 5 is Saturday, 6 is Sunday


from datetime import timedelta

def calculate_duration_excluding_weekends(start_time, end_time):
    """Calculate duration in hours between two datetime objects, excluding weekends."""
    if not isinstance(start_time, datetime) or not isinstance(end_time, datetime):
        raise TypeError("start_time and end_time must be datetime objects")

    if start_time >= end_time:
        return 0

    total_hours = 0
    current_time = start_time

    while current_time < end_time:
        # Check if the current day is a weekday (0-4 corresponds to Monday to Friday)
        if current_time.weekday() < 5:
            # Calculate the end of the current day
            end_of_day = datetime(current_time.year, current_time.month, current_time.day, 23, 59, 59, tzinfo=current_time.tzinfo)
            next_time = min(end_of_day + timedelta(seconds=1), end_time)
            duration = next_time - current_time

            # Convert duration to hours
            hours = duration.total_seconds() / 3600
            total_hours += hours

        # Move to the start of the next day
        current_time = datetime(current_time.year, current_time.month, current_time.day, tzinfo=current_time.tzinfo) + timedelta(days=1)

    return total_hours
import pytz
tz = pytz.UTC

from datetime import datetime
import pytz
def calculate_points_excluding_weekends(start_time, end_time, points_per_hour, start_sprint, end_sprint):
    """Calculate points excluding weekends between two timezone-aware datetimes."""
    if not isinstance(start_time, datetime) or not isinstance(end_time, datetime):
        raise TypeError("start_time and end_time must be datetime objects")

    if start_time.tzinfo is None or end_time.tzinfo is None:
        raise ValueError("start_time and end_time must be timezone-aware datetime objects")

    tz = pytz.UTC  # Or any timezone you need

    # Convert string dates to timezone-aware datetime objects
    if isinstance(start_sprint, str):
        start_sprint = datetime.strptime(start_sprint, "%d/%b/%y").replace(hour=9, minute=0, tzinfo=tz)
    elif start_sprint.tzinfo is None:
        start_sprint = tz.localize(start_sprint.replace(hour=9, minute=0))

    if isinstance(end_sprint, str):
        end_sprint = datetime.strptime(end_sprint, "%d/%b/%y").replace(hour=17, minute=0, tzinfo=tz)
    elif end_sprint.tzinfo is None:
        end_sprint = tz.localize(end_sprint.replace(hour=17, minute=0))

    # Ensure start_time and end_time are within the sprint period
    if start_time < start_sprint:
        start_time = start_sprint
    if end_time > end_sprint:
        end_time = end_sprint

    # Adjust start_time if it is after 17h
    if start_time.hour >= 17:
        start_time = start_time + timedelta(days=1)
        start_time = start_time.replace(hour=9, minute=0)

    # Ensure start_time is before end_time
    if start_time >= end_time:
        return 0

    total_points = 0
    current_time = start_time

    # Assuming calculate_duration_excluding_weekends is a valid function defined elsewhere
    hours_worked = calculate_duration_excluding_weekends(start_time, end_time)
    day_points = hours_worked // 24
    hours_worked = hours_worked % 24

    print("nbre de jour est ", day_points, " nbre d heures est ", hours_worked)
    if hours_worked >= 8:
        day_points += 1
    else:
        day_points += hours_worked * points_per_hour

    return day_points





def calculate_consumed_story_points(issue,start_sprint,end_sprint):
    """Calculate consumed story points for a given issue."""
    status_changes = []

    for history in issue.changelog.histories:
        for item in history.items:
            if item.field == 'status':
                status_changes.append({
                    'from': item.fromString,
                    'to': item.toString,
                    'date': datetime.strptime(history.created, "%Y-%m-%dT%H:%M:%S.%f%z")
                })

    status_changes.sort(key=lambda x: x['date'])

    if not status_changes:
        return 0

    total_points = 0
    points_per_hour = 1/8
    in_progress_start_time = None
    # Convertir start_sprint et end_sprint en offset-aware
    start_sprint = start_sprint.astimezone(pytz.utc)
    end_sprint = end_sprint.astimezone(pytz.utc)

    for change in status_changes:
        from_status = change['from']
        to_status = change['to']
        current_time = change['date']

        if from_status != 'In Progress' and to_status == 'In Progress':
            in_progress_start_time = current_time
            if start_sprint < in_progress_start_time < end_sprint:
                logged=True
            else:
                logged=False


        elif from_status == 'In Progress' and to_status != 'In Progress':
            if start_sprint < in_progress_start_time < end_sprint:
                logged=True
            else:
                logged=False


            if in_progress_start_time:
                total_points += calculate_points_excluding_weekends(in_progress_start_time, current_time, points_per_hour,start_sprint,end_sprint)
                in_progress_start_time = None


    if in_progress_start_time:
        current_time = datetime.now(pytz.utc)
        total_points += calculate_points_excluding_weekends(in_progress_start_time, current_time, points_per_hour,start_sprint,end_sprint)

    time_estimate = issue.fields.timeoriginalestimate
    time_spent = issue.fields.timespent

    if time_estimate:
        time_estimate_hours = time_estimate / 3600  # Convertir en heures
        print(f"Estimated Time: {time_estimate_hours} hours")
    else:
        print("Estimated Time: Not provided")

    if time_spent:
        time_spent_hours = time_spent / 3600  # Convertir en heures
        if logged:
            total_points=time_spent_hours/8
            print("Use Logged")

        print(f"Logged Time: {time_spent_hours} hours")
    else:
        print("Logged Time: Not provided")

    return total_points


def generate_velocity_chart(sprint_data, filename):
    """Generate and save a velocity chart based on sprint data."""
    sprint_names = [data['Sprint Name'] for data in sprint_data]
    total_points = [data['Total Points'] for data in sprint_data]

    plt.figure(figsize=(10, 6))
    plt.plot(sprint_names, total_points, marker='o', linestyle='-', color='b')
    plt.xlabel('Sprint')
    plt.ylabel('Total Points')
    plt.title('Velocity Chart')
    plt.xticks(rotation=45)
    plt.grid(True)

    plt.tight_layout()
    plt.savefig(filename)
    plt.close()


def print_issue_details(issues, start_sprint, end_sprint):
    """Print details of the issues."""
    for issue in issues:
        print(f"Issue Key: {issue.key}")
        print(f"Summary: {issue.fields.summary}")

        reporter = issue.fields.reporter
        if reporter:
            print(f"Reporter: {reporter.displayName}")
        else:
            print("Reporter: Unassigned")

        assignee = issue.fields.assignee
        if assignee:
            print(f"Assignee: {assignee.displayName}")
        else:
            print("Assignee: Unassigned")

        print(f"Status: {issue.fields.status.name}")
        print(f"Created: {issue.fields.created}")
        print(f"Updated: {issue.fields.updated}")

        # Consumed Story Points
        print(f"Consumed Story Points: {calculate_consumed_story_points(issue, start_sprint, end_sprint)}")

        # Suivi du temps (temps estimé et temps enregistré)
        time_estimate = issue.fields.timeoriginalestimate
        time_spent = issue.fields.timespent

        if time_estimate:
            time_estimate_hours = time_estimate / 3600  # Convertir en heures
            print(f"Estimated Time: {time_estimate_hours} hours")
        else:
            print("Estimated Time: Not provided")

        if time_spent:
            time_spent_hours = time_spent / 3600  # Convertir en heures
            print(f"Logged Time: {time_spent_hours} hours")
        else:
            print("Logged Time: Not provided")

        # Affichage de l'historique des transitions
        print("\nHistory:\n")
        for history in issue.changelog.histories:
            transition_items = [item for item in history.items if item.field == 'status']
            if transition_items:
                print(f"Date: {history.created}")
                for item in transition_items:
                    print(f"Changed {item.field.capitalize()}: From '{item.fromString}' To '{item.toString}'")
                print("-" * 20)

        print("\n" + "=" * 20 + "\n")

def get_issue_by_key(issue_key):
    """Retrieve issue details by issue key."""
    issue = jira.issue(issue_key, expand='changelog')
    return issue

def excel_serial_to_date(serial_date):
    """Convert Excel serial date to pandas Timestamp."""
    return pd.to_datetime('1899-12-30') + pd.to_timedelta(serial_date, 'D')

app = Flask(__name__)
CORS(app)  # Enable CORS for the entire application


def get_db_connection():
    try:
        connection = mysql.connector.connect(
            host=serverdb,
            port='3306',
            database='actiabackImg',
            user='root',
            password=''
        )
        return connection
    except Error as e:
        print(f"Error connecting to MySQL: {e}")
        return None

def get_sprint_from_key(file, key_value):
    # Read the Excel file without headers to find the header row
    data = pd.read_excel(file, header=None)

    # Find the header row by checking for the expected header pattern
    for i, row in data.iterrows():
        if 'Key' in row.values and 'Issue Type' in row.values:  # Adjust conditions based on unique header columns
            header_row_index = i
            break

    # Read the Excel file again with the detected header row
    data = pd.read_excel(file, header=header_row_index)

    # Trim any leading or trailing spaces from column names
    data.columns = data.columns.str.strip()

    # Check if 'Key' and 'Status' columns are present
    if 'Key' not in data.columns or 'Sprint' not in data.columns:
        raise ValueError("The required columns 'Key' and 'Status' are not present in the Excel file.")

    # Find the row where the 'Key' matches the provided key_value
    issue_row = data[data['Key'] == key_value]

    # Check if the issue was found
    if issue_row.empty:
        return f"No issue found with Key: {key_value}"

    # Return the status of the issue
    return issue_row['Sprint'].values[0]

def get_status_from_key(file, key_value):
    # Read the Excel file without headers to find the header row
    data = pd.read_excel(file, header=None)

    # Find the header row by checking for the expected header pattern
    for i, row in data.iterrows():
        if 'Key' in row.values and 'Issue Type' in row.values:  # Adjust conditions based on unique header columns
            header_row_index = i
            break

    # Read the Excel file again with the detected header row
    data = pd.read_excel(file, header=header_row_index)

    # Trim any leading or trailing spaces from column names
    data.columns = data.columns.str.strip()

    # Check if 'Key' and 'Status' columns are present
    if 'Key' not in data.columns or 'Status' not in data.columns:
        raise ValueError("The required columns 'Key' and 'Status' are not present in the Excel file.")

    # Find the row where the 'Key' matches the provided key_value
    issue_row = data[data['Key'] == key_value]

    # Check if the issue was found
    if issue_row.empty:
        return f"No issue found with Key: {key_value}"

    # Return the status of the issue
    return issue_row['Status'].values[0]
def arrondir_a_025(valeur):
    return round(valeur * 4) / 4

@app.route('/submit-form/<int:idsprint>', methods=['POST'])
def submit_form(idsprint):

    print(idsprint)
    from datetime import datetime

    try:
        # Extract form data
        sprint_name = request.form.get('sprintName')
        raw_start_date = request.form.get('startDate')
        raw_end_date = request.form.get('endDate')

        # Parse only the first part of the date string
        parsed_start_date = datetime.strptime(raw_start_date[:24], '%a %b %d %Y %H:%M:%S')
        parsed_end_date = datetime.strptime(raw_end_date[:24], '%a %b %d %Y %H:%M:%S')

        # Reformat the dates to match '%d/%b/%y'
        start_date = parsed_start_date.strftime('%d/%b/%y')
        end_date = parsed_end_date.strftime('%d/%b/%y')
        carry_forward_sp = request.form.get('carryForwardSP')

        # Extract file data
        file = request.files.get('file')

        file_data = None
        if file:
            file_data = file.read()
            # Optionally print file size or type
            print(f"File size: {len(file_data)} bytes")
            print(f"File type: {file.mimetype}")




        # Print extracted sprint details
        print(f"Sprint Name: {sprint_name}")
        print(f"Start Date: {start_date}")
        print(f"End Date: {end_date}")
        print(f"Carry Forward SP: {carry_forward_sp}")

        if file:
            print("File received")
        else:
            print("No file received")



        from datetime import datetime, timedelta

        # Importation du capture d'ecran



        def get_valid_date(start_date, end_date):
            try:
                start_sprint = datetime.strptime(start_date, "%d/%b/%y").replace(hour=9, minute=0)
                end_sprint = datetime.strptime(end_date, "%d/%b/%y").replace(hour=17, minute=0)
                if start_sprint <= end_sprint:
                    return start_sprint, end_sprint
                else:
                    raise ValueError("Start date must be earlier or equal to end date.")
            except ValueError as e:
                print("Error parsing dates:", e)
                raise

        # Use the dates passed as parameters
        start_sprint, end_sprint = get_valid_date(start_date, end_date)
        print("Date debut sprint:", start_sprint)
        print("Date fin sprint:", end_sprint)
        # Specifiez le chemin complet du fichier sur votre bureau
        # ***************
        #chemin_fichier = r"C:\Users\azizs\Desktop\dataexel\MATRIX_24.2-1.xlsx"
        # Importez le fichier en utilisant pandas


        data = pd.read_excel(file, header=None)
        # Find the header row by checking for the expected header pattern
        for i, row in data.iterrows():
            if 'Key' in row.values and 'Issue Type' in row.values:  # Adjust conditions based on unique header columns
                header_row_index = i
                break

        # Read the Excel file again with the detected header row
        data = pd.read_excel(file, header=header_row_index)


        # Drop rows where the "Key" column is empty
        data = data.dropna(subset=["Key"])

        # Convert serial dates to human-readable dates
        date_columns = ['Start Date', 'End Date']  # Replace with your actual date column names
        for col in date_columns:
            if col in data.columns:
                data[col] = data[col].apply(excel_serial_to_date)

        # Initialize an empty list to store results
        results = []
        data = data.dropna(subset=["Key"])

        # Iterate through each issue key from the Excel file
        for issues_key in data["Key"]:

            issue = get_issue_by_key(issues_key)
            print(f"Issue Key: {issue.key}")
            print(f"Summary: {issue.fields.summary}")

            # Handle reporter field
            reporter = issue.fields.reporter
            if reporter:
                print(f"Reporter: {reporter.displayName}")
            else:
                print("Reporter: Unassigned")

            # Handle assignee field
            assignee = issue.fields.assignee
            if assignee:
                print(f"Assignee: {assignee.displayName}")
            else:
                print("Assignee: Unassigned")

            print(f"Status: {issue.fields.status.name}")
            print(f"Created: {issue.fields.created}")
            print(f"Updated: {issue.fields.updated}")
            print(f"Consumed Story Points: {calculate_consumed_story_points(issue,start_sprint,end_sprint)}")
            print("\nHistory:\n")
            time_estimate = issue.fields.timeoriginalestimate
            time_spent = issue.fields.timespent

            if time_estimate:
                time_estimate_hours = time_estimate / 3600  # Convertir en heures
                print(f"Estimated Time: {time_estimate_hours} hours")
            else:
                print("Estimated Time: Not provided")

            if time_spent:
                time_spent_hours = time_spent / 3600  # Convertir en heures
                print(f"Logged Time: {time_spent_hours} hours")
            else:
                print("Logged Time: Not provided")

            # Display the issue's history (changelog)
            for history in issue.changelog.histories:
                transition_items = [item for item in history.items if item.field == 'status']
                if transition_items:
                    print(f"Date: {history.created}")
                    for item in transition_items:
                        print(f"Changed {item.field.capitalize()}: From '{item.fromString}' To '{item.toString}'")
                    print("-" * 20)
            print("\n" + "=" * 20 + "\n")
            if issue:
                consumed_story_points = calculate_consumed_story_points(issue,start_sprint,end_sprint)
                latest_status = get_status_from_key(file,issues_key)
                components = ', '.join(comp.name for comp in issue.fields.components) if hasattr(issue.fields,'components') else 'None'
               # story_points = issue.fields.customfield_10016 if hasattr(issue.fields, 'customfield_10016') else 'N/A'
                sprint = get_sprint_from_key(file,issues_key)

                history = [
                    {
                        'date': datetime.strptime(history.created, "%Y-%m-%dT%H:%M:%S.%f%z"),
                        'from': item.fromString,
                        'to': item.toString
                    }
                    for history in issue.changelog.histories
                    for item in history.items
                    if item.field == 'status'
                ]
                results.append({
                    'Issue Key': issue.key,
                    'Issue Type': issue.fields.issuetype.name,

                    'Summary': issue.fields.summary,
                    'Reporter': issue.fields.reporter.displayName if issue.fields.reporter else 'Unassigned',
                    'Assignee': issue.fields.assignee.displayName if issue.fields.assignee else 'Unassigned',
                    'Status': latest_status,
                    'Created': datetime.strptime(issue.fields.created, "%Y-%m-%dT%H:%M:%S.%f%z").strftime(
                        '%d-%b-%Y %H:%M'),
                    'Updated': datetime.strptime(issue.fields.updated, "%Y-%m-%dT%H:%M:%S.%f%z").strftime(
                        '%d-%b-%Y %H:%M'),
                    'Component/s': components,
                    #'Story Points': story_points,
                    'Sprint': sprint,

                    'consumed Story Points': arrondir_a_025(consumed_story_points),

                })

        results = [result for result in results if result['Issue Type'] != 'Epic' and result['consumed Story Points'] != 0]

        # Export to Excel
        output_file = 'jira_issues_with_history.xlsx'
        df = pd.DataFrame(results)
        df.to_excel(output_file, index=False)

        # Use 'Story Points' column instead of generating random values
        #df['Treated task'] = np.random.randint(0, 5, df.shape[0])

        # Removing leading/trailing spaces from 'Status' column
        df['Status'] = df['Status'].str.strip()

        # Grouping by 'Status' and summing 'Story Points'
        sum_by_status = df.groupby('Status')['consumed Story Points'].sum().reset_index()
        sum_by_status = df.groupby('Status')['consumed Story Points'].sum().reset_index()
        #sum_treated_task = df.groupby(['Issue Type', 'Status'])['Treated task'].sum().reset_index()
        #count_status = df.groupby('Issue Type')['Status'].value_counts().unstack().reset_index()

        # Filtering statuses for 'Integrated' and 'Closed'
        filtered_status = sum_by_status[sum_by_status['Status'].isin(['Integrated', 'Closed'])]

        # Summing the 'Story Points' for the filtered statuses
        total_sum = filtered_status['consumed Story Points'].sum()

        # Convert the result to JSON
        json_result_sum_by_status = sum_by_status.to_json(orient='records')
        print(json_result_sum_by_status)

        # Counting issues by 'Issue Type'
        issue_counts = df["Issue Type"].value_counts().reset_index()
        issue_counts.columns = ['Issue Type', 'Count']

        # Convert the issue counts DataFrame to JSON
        json_result_histogram = issue_counts.to_json(orient='records')
        print(json_result_histogram)

        # Identifying treated issues with 'Integrated' or 'Closed' status
        df['Treated issues (Integrated or Closed)'] = df['Status'].apply(
            lambda status: status if status in ['Closed', 'Integrated'] else np.nan
        )

        treated_counts = df.groupby('Issue Type')['Treated issues (Integrated or Closed)'].apply(
            lambda x: x.notna().sum()
        ).reset_index()

        treated_counts.columns = ['Issue Type', 'Treated issues (Integrated or Closed)']

        # Merging the counts with treated issues
        issue_df = pd.merge(issue_counts, treated_counts, on='Issue Type', how='left')

        # Convert the merged DataFrame to JSON
        json_result_grouped_bar_chart = issue_df.to_json(orient='records')
        print(json_result_grouped_bar_chart)

        import mysql.connector
        from mysql.connector import Error

        try:
            # Connect to MySQL
            connection = mysql.connector.connect(
                host=serverdb,
                port='3306',
                database='actiabackImg',
                user='root',
                password=''
            )

            if connection.is_connected():
                cursor = connection.cursor()

                # SQL statement to create the 'files' table
                create_table_query = """
                CREATE TABLE files (
                    id INT AUTO_INCREMENT PRIMARY KEY,
                    file_name VARCHAR(255),
                    file_data LONGBLOB
                );
                """

                # SQL statement to alter the 'performanceJira' table
                alter_table_query = """
                ALTER TABLE performanceJira
                ADD COLUMN file_data LONGBLOB;
                """

                # Execute the queries
                cursor.execute(create_table_query)
                print("Table 'files' created successfully.")

                cursor.execute(alter_table_query)
                print("Column 'file_data' added to 'performanceJira' table successfully.")

        except Error as e:
            print(f"Error: {e}")

        finally:
            if connection.is_connected():
                cursor.close()
                connection.close()
                print("MySQL connection closed")

        import mysql.connector
        from mysql.connector import Error

        try:
            # Connect to MySQL
            connection = mysql.connector.connect(
                host=serverdb,
                port='3306',
                database='actiabackImg',
                user='root',
                password=''
            )

            if connection.is_connected():
                cursor = connection.cursor()

                # SQL statement to alter the table and add the 'file_data' column
                alter_table_query = """
                ALTER TABLE performanceJira
                ADD COLUMN file_data LONGBLOB;
                """

                # Execute the query
                cursor.execute(alter_table_query)
                connection.commit()
                print("Column 'file_data' added to 'performanceJira' table successfully.")

        except Error as e:
            print(f"Error: {e}")

        finally:
            if connection.is_connected():
                cursor.close()
                connection.close()
                print("MySQL connection closed")

        import mysql.connector
        from mysql.connector import Error

        # File path

        # Example data


        try:
            # Connect to MySQL
            connection = mysql.connector.connect(
                host=serverdb,
                port='3306',
                database='actiabackImg',
                user='root',
                password=''
            )

            if connection.is_connected():
                cursor = connection.cursor()

                # Create a BytesIO object to hold the Excel data in memory
                output_file = io.BytesIO()

                # Write the DataFrame to this in-memory file
                df = pd.DataFrame(results)
                df.to_excel(output_file, index=False)

                # Get the in-memory data
                file_data = output_file.getvalue()
                # Insert data query
                insert_data_query = """
                           INSERT INTO performanceJira (start_date, end_date, sprint_name, grouped_bar_chart, pie_chart, histogram, file_data,idsprint)
                           VALUES (%s, %s, %s, %s, %s, %s, %s,%s)
                           """

                cursor.execute(insert_data_query, (
                    datetime.strptime(raw_start_date[:24], '%a %b %d %Y %H:%M:%S').date(),
                    datetime.strptime(raw_end_date[:24], '%a %b %d %Y %H:%M:%S').date(),
                    sprint_name,
                    json.dumps(json_result_grouped_bar_chart),
                    json.dumps(json_result_sum_by_status),
                    json.dumps(json_result_histogram),
                    file_data,
                    idsprint
                ))
                connection.commit()
                print("Data and file inserted successfully.")

        except Error as e:
            print(f"Error: {e}")

        return jsonify({"message": "Form data received successfully"})
    except Exception as e:
        print(f"Error processing form data: {e}")
        return jsonify({"error": "Failed to process form data"}), 500
from flask import jsonify


@app.route('/data/<int:idsprint>', methods=['GET'])
def get_data_by_idsprint(idsprint):
    connection = get_db_connection()
    if not connection:
        return jsonify({"error": "Failed to connect to the database"}), 500

    try:
        cursor = connection.cursor(dictionary=True)
        query = "SELECT * FROM performanceJira WHERE idsprint = %s"
        cursor.execute(query, (idsprint,))
        result = cursor.fetchall()

        # Deserialize JSON fields and encode binary data
        for row in result:
            if row['grouped_bar_chart']:
                row['grouped_bar_chart'] = json.loads(row['grouped_bar_chart'])
            if row['pie_chart']:
                row['pie_chart'] = json.loads(row['pie_chart'])
            if row['histogram']:
                row['histogram'] = json.loads(row['histogram'])
            # Encode binary file data to base64
            if row['file_data']:
                row['file_data'] = base64.b64encode(row['file_data']).decode('utf-8')

        return jsonify(result)
    except Error as e:
        print(f"Error fetching data from MySQL: {e}")
        return jsonify({"error": "Failed to fetch data from the database"}), 500
    finally:
        if connection.is_connected():
            cursor.close()
            connection.close()


@app.route('/download-xlsx/<int:file_id>', methods=['GET'])
def download_xlsx(file_id):
    try:
        # Connect to MySQL
        connection = mysql.connector.connect(
            host=serverdb,
            port='3306',
            database='actiabackImg',
            user='root',
            password=''
        )

        if connection.is_connected():
            cursor = connection.cursor()
            query = "SELECT file_data FROM performanceJira WHERE id = %s"
            cursor.execute(query, (file_id,))
            result = cursor.fetchone()

            if result:
                file_data = result[0]

                # Create an in-memory file object
                file_stream = io.BytesIO(file_data)

                # Send the file as a response
                return send_file(
                    file_stream,
                    as_attachment=True,
                    download_name=f'performanceJira_{file_id}.xlsx',
                    mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
                )
            else:
                return jsonify({"error": "File not found"}), 404

    except Error as e:
        return jsonify({"error": str(e)}), 500


# Charger la présentation existante
prs = Presentation("MATRIX-Presentation.pptx")

# Obtenir la date actuelle avec le mois en majuscule (par exemple, "SEPTEMBER2024")
date_actuelle = date_actuelle = datetime.now().strftime("%Y/%m/%d")

date_actuelle2 = datetime.now().strftime("%B%Y").upper()  # Cela donne "SEPTEMBER2024", "NOVEMBER2024", etc.


def remplacer_textes(slide,sprintname):
    # Remplacer la dernière occurrence de "PI" par "PMT"
    if "PI" in sprintname:
        sprintname = sprintname.rsplit("PI", 1)  # Diviser la chaîne en deux parties
        sprintname = sprintname[0] + "PMT" + sprintname[1]  # Reconstituer la chaîne avec "PMT"

    # Afficher le résultat
    print(sprintname)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            full_text = "".join([run.text for run in paragraph.runs])
            # Remplacer la date actuelle
            if "2024/09/05" in full_text:
                for run in paragraph.runs:
                    run.text = run.text.replace("2024/09/05", date_actuelle)
            # Remplacer le texte spécifique avec un match plus flexible
            if "MATRIX - sEPTEMBER2024 – PMt24.3-3" in full_text:
                new_text = f"MATRIX - {date_actuelle2} – "+sprintname
                for run in paragraph.runs:
                    run.text = run.text.replace("MATRIX - sEPTEMBER2024 – PMt24.3-3", new_text)


            if "MATRIX - SEPTEMBER2024- Pmt24.3-3" in full_text:
                new_text = f"MATRIX - {date_actuelle2} – "+sprintname
                for run in paragraph.runs:
                    run.text = run.text.replace("MATRIX - SEPTEMBER2024- Pmt24.3-3", new_text)

            if "PMT24.3-3 September2024" in full_text:
                new_text2 = f"{sprintname} {date_actuelle2}"
                for run in paragraph.runs:
                    run.text = run.text.replace("PMT24.3-3 September2024", new_text2)

# Fonction pour remplacer l'image sur une slide donnée

def remplacer_image(slide, image_path, image_size=(5,4)):
    slide_width = prs.slide_width  # Largeur de la diapositive
    slide_height = prs.slide_height  # Hauteur de la diapositive

    # Calculer les positions pour centrer l'image
    left = (slide_width - Inches(image_size[0])) / 2
    top = (slide_height - Inches(image_size[1])) / 2

    # Supprimer les anciennes images de la diapositive
    for shape in slide.shapes:
        if shape.shape_type == 13:  # 13 correspond à un type "image"
            sp = shape
            slide.shapes._spTree.remove(sp._element)

    # Ajouter une nouvelle image centrée
    slide.shapes.add_picture(image_path, left, top, Inches(image_size[0]), Inches(image_size[1]))

sum_by_status=None
@app.route('/download-presentation/<int:file_id>', methods=['POST'])
def download_presentation(file_id):
    carry_forward_sp = request.form.get('capacity')
    sprintName = request.form.get('sprintName')
    file = request.files.get('file')
    file_vi = None

    fileexel = request.files.get('fileexel')

    import mysql.connector
    import io
    from openpyxl import load_workbook

    try:
        # Connexion à MySQL
        connection = mysql.connector.connect(
            host=serverdb,
            port='3306',
            database='actiabackImg',
            user='root',
            password=''
        )

        if connection.is_connected():
            cursor = connection.cursor()

            if fileexel:
                print("Fichier Excel reçu")

                # Lire les données du fichier
                file_data = fileexel.read()

                # Charger le fichier Excel avec openpyxl
                workbookexel = load_workbook(filename=io.BytesIO(file_data))

                sheet = workbookexel.active  # Ou sélectionner une autre feuille si nécessaire

                # Extraire les en-têtes (première ligne)
                headers = [cell.value for cell in sheet[1]]

                # Vérifier la présence de la colonne 'consumed Story Points'
                if 'consumed Story Points' in headers:
                    print("La colonne 'consumed Story Points' est présente.")

                    # Requête pour mettre à jour les données dans MySQL
                    update_data_query = """
                        UPDATE performanceJira
                        SET file_data = %s
                        WHERE id = %s
                    """

                    # Exécution de la requête avec les données correspondantes
                    cursor.execute(update_data_query, (
                        file_data,  # Utiliser le fichier Excel uploadé
                        file_id  # Assurez-vous que file_id est défini avant
                    ))

                    # Commit des changements pour valider la mise à jour
                    connection.commit()

                    print(f"Les données ont été mises à jour avec succès pour l'ID {file_id}.")
                else:
                    print("La colonne 'consumed Story Points' est absente. Aucune mise à jour effectuée.")


        if connection.is_connected():
            cursor = connection.cursor()
            query = """
                SELECT file_data
                FROM performanceJira 
                WHERE id = %s
            """
            cursor.execute(query, (file_id,))
            result = cursor.fetchone()

            if result:
                file_data = result[0]


                # Créer un fichier en mémoire à partir des données Excel
                file_stream = io.BytesIO(file_data)

                # Lire le fichier Excel en utilisant pandas
                try:
                    excel_data = pd.read_excel(file_stream)

                    # Vérifier si la colonne 'Consumed Story Points' existe
                    if 'consumed Story Points' not in excel_data.columns:
                        return jsonify(
                            {"error": "La colonne 'Consumed Story Points' n'existe pas dans le fichier Excel."}), 400

                    # Calculer la somme de 'Consumed Story Points'
                    sum_by_status = excel_data.groupby('Status')['consumed Story Points'].sum().reset_index()
                    filtered_status = sum_by_status[sum_by_status['Status'].isin(['Integrated', 'Closed'])]
                    total_sum = filtered_status['consumed Story Points'].sum()

                    excel_data['Status'] = excel_data['Status'].str.strip()

                    # Grouping by 'Status' and summing 'Story Points'
                    sum_by_status = excel_data.groupby('Status')['consumed Story Points'].sum().reset_index()
                    sum_by_status = excel_data.groupby('Status')['consumed Story Points'].sum().reset_index()
                    # sum_treated_task = df.groupby(['Issue Type', 'Status'])['Treated task'].sum().reset_index()
                    # count_status = df.groupby('Issue Type')['Status'].value_counts().unstack().reset_index()

                    # Filtering statuses for 'Integrated' and 'Closed'
                    filtered_status = sum_by_status[sum_by_status['Status'].isin(['Integrated', 'Closed'])]

                    # Summing the 'Story Points' for the filtered statuses
                    total_sum = filtered_status['consumed Story Points'].sum()

                    # Convert the result to JSON
                    json_result_sum_by_status = sum_by_status.to_json(orient='records')
                    print(json_result_sum_by_status)

                    # Counting issues by 'Issue Type'
                    issue_counts = excel_data["Issue Type"].value_counts().reset_index()
                    issue_counts.columns = ['Issue Type', 'Count']

                    # Convert the issue counts DataFrame to JSON
                    json_result_histogram = issue_counts.to_json(orient='records')
                    print(json_result_histogram)



                    # Identifying treated issues with 'Integrated' or 'Closed' status
                    excel_data['Treated issues (Integrated or Closed)'] = excel_data['Status'].apply(
                        lambda status: status if status in ['Closed', 'Integrated'] else np.nan
                    )

                    treated_counts = excel_data.groupby('Issue Type')['Treated issues (Integrated or Closed)'].apply(
                        lambda x: x.notna().sum()
                    ).reset_index()

                    treated_counts.columns = ['Issue Type', 'Treated issues (Integrated or Closed)']

                    # Merging the counts with treated issues
                    issue_df = pd.merge(issue_counts, treated_counts, on='Issue Type', how='left')

                    # Convert the merged DataFrame to JSON
                    json_result_grouped_bar_chart = issue_df.to_json(orient='records')
                    print(json_result_grouped_bar_chart)
                    data_grouped_bar_chart = json.loads(json_result_grouped_bar_chart)

                    issue_types = [item['Issue Type'] for item in data_grouped_bar_chart]
                    all_issues = [item['Count'] for item in data_grouped_bar_chart]
                    treated_issues = [item['Treated issues (Integrated or Closed)'] for item in data_grouped_bar_chart]
                    print(total_sum)
                    issue_types = [item['Issue Type'] for item in data_grouped_bar_chart]
                    all_issues = [item['Count'] for item in data_grouped_bar_chart]
                    treated_issues = [item['Treated issues (Integrated or Closed)'] for item in data_grouped_bar_chart]

                    update_data_query2 = """
                        UPDATE performanceJira
                        SET grouped_bar_chart = %s, pie_chart = %s, histogram = %s
                        WHERE id = %s
                    """

                    # Exécution de la requête avec les données correspondantes et l'ID de la ligne à mettre à jour
                    cursor.execute(update_data_query2, (
                        json.dumps(json_result_grouped_bar_chart),  # Sérialisation en JSON du graphe groupé
                        json.dumps(json_result_sum_by_status),  # Sérialisation en JSON du graphique en camembert
                        json.dumps(json_result_histogram),  # Sérialisation en JSON de l'histogramme
                        file_id  # Assurez-vous que file_id (l'ID de l'enregistrement) est défini
                    ))

                    # Commit des changements pour appliquer la mise à jour
                    connection.commit()

                    print("Les données ont été mises à jour avec succès pour l'ID", file_id)

                    # Create a figure and axis
                    fig, ax = plt.subplots()

                    # Define the position of the bars on the y-axis
                    index = np.arange(len(issue_types))

                    # Define the bar width
                    bar_width = 0.35

                    # Plot the bars: One for All Issues and one for Treated Issues
                    bar1 = ax.barh(index - bar_width / 2, all_issues, bar_width, label='All Issues', color='lightcoral')
                    bar2 = ax.barh(index + bar_width / 2, treated_issues, bar_width,
                                   label='Treated Issues (Integrated or Closed)',
                                   color='lightgreen')

                    # Add chart title and labels
                    ax.set_title(sprintName+' MATRIX Status By Issue Type')
                    ax.set_xlabel('Number of Issues')
                    ax.set_yticks(index)
                    ax.set_yticklabels(issue_types)
                    ax.invert_yaxis()  # Invert y-axis to match the order of issue types (optional)

                    # Set x-ticks to display 0, 2, 4, 6, 8, etc.
                    ax.set_xticks(np.arange(0, max(all_issues)+4, 2))  # Customize the range and step (0 to 20, step 2)

                    # Add grid lines for readability
                    ax.grid(True, which='both', axis='x', linestyle='--', linewidth=0.5)

                    # Add a legend
                    ax.legend()

                    # Save the plot to a file
                    plt.tight_layout()
                    plt.savefig("image1.png")
                    plt.savefig("image2.png")
                    print(f"Chart saved as image1.png")
                    # Handle the charts (grouped_bar_chart and pie_chart) if needed here
                    # For example, you can convert them to JSON if they are stored in a compatible format
                    # or return them as part of the response

                except Exception as e:
                    return jsonify({"error": f"Erreur lors de la lecture du fichier Excel: {str(e)}"}), 500
            else:
                return jsonify({"error": "Le fichier avec l'ID spécifié n'existe pas."}), 404


    except mysql.connector.Error as e:
        return jsonify({"error": f"Erreur de connexion à la base de données: {str(e)}"}), 500

    except Exception as e:
        return jsonify({"error": f"Erreur inattendue: {str(e)}"}), 500

        # Initialize the new structure with headers

    # Initialize the workbook and sheet
    from openpyxl import Workbook
    from openpyxl.chart import PieChart3D, Reference
    from openpyxl.chart.label import DataLabelList

    # Créer un ordre personnalisé pour les statuts
    custom_order = [
        "test",
        "In Progress",
        "Implemented",
        "Integrated",
        "Closed",
        "In Review",
        "Blocked"
    ]

    # Ajouter une valeur spéciale pour gérer les autres statuts (pas dans custom_order)
    sum_by_status['sort_key'] = sum_by_status['Status'].apply(
        lambda x: custom_order.index(x) if x in custom_order else len(custom_order)
    )

    # Trier le DataFrame en fonction de cet ordre
    sum_by_status = sum_by_status.sort_values(by='sort_key')

    # Supprimer la colonne de clé de tri après le tri
    sum_by_status = sum_by_status.drop(columns=['sort_key'])

    # Créer le fichier Excel et le graphique
    workbook = Workbook()
    sheet = workbook.active

    data = [['Status', 'Consumed Story Points']]
    data.append(['test',10])
    # Remplir la liste de données
    for index, row in sum_by_status.iterrows():
        status = row['Status']
        consumed_story_points = row['consumed Story Points']
        data.append([status, consumed_story_points])

    # Écrire les données dans la feuille
    for row in data:
        sheet.append(row)

    # Créer un diagramme circulaire 3D
    pie = PieChart3D()

    # Définir les plages de données pour les valeurs et les étiquettes
    labels = Reference(sheet, min_col=1, min_row=3, max_row=len(data) + 1)
    data_values = Reference(sheet, min_col=2, min_row=2, max_row=len(data) + 1)

    # Ajouter les données et les catégories au graphique
    pie.add_data(data_values, titles_from_data=True)
    pie.set_categories(labels)
    pie.title = sprintName + " MATRIX Story Points Status"

    # Ajouter des étiquettes de données
    pie.dLbls = DataLabelList()
    pie.dLbls.showVal = True
    pie.dLbls.showSerName = False
    pie.dLbls.showPercent = False
    pie.dLbls.showCatName = True
    pie.dLbls.dLblPos = 'bestFit'  # Positionner les étiquettes à l'extérieur du graphique
    pie.legend.position = 'b'

    # Définir le style et ajuster la taille
    pie.style = 26
    pie.height = 10
    pie.width = 18

    # Appliquer des couleurs spécifiques aux segments du diagramme
    color_mapping = {
        "Implemented": "99FF99",  # Vert clair
        "Closed": "336633",  # Vert foncé
        "Integrated": "339933",  # Vert moyen
        "Blocked": "999999",  # Gris
        "In Progress": "FFFF66",  # Jaune
        "Backlog Refinement": "FFCC99",  # Beige / Orange clair
        "In Review": "FF9966",  # Orange clair
        "Ready": "66CCFF"  # Bleu clair
    }

    # Appliquer les couleurs aux points de données
    for i, status in enumerate(sum_by_status['Status']):
        point = DataPoint(i)
        point.graphicalProperties.solidFill = color_mapping.get(status, "CCCCCC")  # Couleur par défaut si non trouvé
        pie.series[0].data_points.append(point)

    # Positionner le graphique dans la feuille Excel
    sheet.add_chart(pie, "E2")

    # Enregistrer le fichier Excel avec le graphique
    workbook.save("exelPie.xlsx")

    print(sprintName)
    # Extract file data
    file = request.files.get('file')
    file_vi = None
    if file:
        file_vi = file.read()
        # Optionally print file size or type
        print(f"File size: {len(file_vi)} bytes")
        print(f"File type: {file.mimetype}")

    # Load the Excel workbook from binary data
    wb = load_workbook(filename=io.BytesIO(file_vi))

    # Charger le fichier Excel

    ws = wb.active

    # Trouver la dernière colonne
    max_col = ws.max_column  # Index de la dernière colonne
    new_col_index = max_col + 1  # Index de la nouvelle colonne

    # Dupliquer la dernière colonne dans la nouvelle colonne
    for row in range(1, ws.max_row + 1):  # Inclure toutes les lignes (y compris l'en-tête)
        old_cell = ws.cell(row=row, column=max_col)
        new_cell = ws.cell(row=row, column=new_col_index, value=old_cell.value)

        # Copier les propriétés de style individuellement (remplissage, police, bordure, alignement)
        if old_cell.fill:
            new_cell.fill = PatternFill(
                start_color=old_cell.fill.start_color,
                end_color=old_cell.fill.end_color,
                fill_type=old_cell.fill.fill_type
            )

        if old_cell.font:
            new_cell.font = Font(
                name=old_cell.font.name,
                size=old_cell.font.size,
                bold=old_cell.font.bold,
                italic=old_cell.font.italic,
                vertAlign=old_cell.font.vertAlign,
                underline=old_cell.font.underline,
                strike=old_cell.font.strike,
                color=old_cell.font.color
            )

        if old_cell.border:
            new_cell.border = Border(
                left=old_cell.border.left,
                right=old_cell.border.right,
                top=old_cell.border.top,
                bottom=old_cell.border.bottom
            )

        if old_cell.alignment:
            new_cell.alignment = Alignment(
                horizontal=old_cell.alignment.horizontal,
                vertical=old_cell.alignment.vertical,
                text_rotation=old_cell.alignment.text_rotation,
                wrap_text=old_cell.alignment.wrap_text,
                shrink_to_fit=old_cell.alignment.shrink_to_fit,
                indent=old_cell.alignment.indent
            )

    # Maintenant modifier les valeurs dans la nouvelle colonne
    new_col_values = [sprintName, int(carry_forward_sp), total_sum]  # Valeurs à mettre dans la nouvelle colonne
    for row, value in enumerate(new_col_values, start=1):  # Commencer à la première ligne
        ws.cell(row=row, column=new_col_index, value=value)

    # Sauvegarder le fichier modifié


    excel_output = io.BytesIO()
    wb.save(excel_output)
    excel_output.seek(0)
    df1 = pd.read_excel(excel_output)

    # Check the structure of the data
    print(df1.head())  # To make sure we have the correct structure

    # Transpose the DataFrame to switch rows with columns
    df1 = df1.transpose()

    # Set the first row (which contains 'Capacity' and 'Velocity') as the new column names
    df1.columns = df1.iloc[0]

    # Drop the first row since it's now the header
    df1 = df1.drop(df1.index[0])

    # Convert the values to numeric (if necessary)
    df1['Capacity'] = pd.to_numeric(df1['Capacity'], errors='coerce')
    df1['Velocity'] = pd.to_numeric(df1['Velocity'], errors='coerce')

    # Plot the graph
    plt.figure(figsize=(10, 6))
    plt.plot(df1.index, df1['Capacity'], marker='o', label='Capacity', color='blue', linewidth=2)
    plt.plot(df1.index, df1['Velocity'], marker='o', label='Velocity', color='yellow', linewidth=2)

    # Adding titles and labels
    plt.title('MATRIX Velocity In PIs')
    plt.xlabel('PIs')

    # Display legend
    plt.legend()

    # Display grid
    plt.grid(True)

    # Rotate x-axis labels for better readability
    plt.xticks(rotation=90)

    # Add value labels above each point for Capacity using iloc for safe positional indexing
    for i in range(len(df1)):
        plt.text(df1.index[i], df1['Capacity'].iloc[i] + 1, str(df1['Capacity'].iloc[i]),
                 ha='center', fontsize=9, color='blue')

    # Add value labels above each point for Velocity using iloc for safe positional indexing
    for i in range(len(df1)):
        plt.text(df1.index[i], df1['Velocity'].iloc[i] + 1, str(df1['Velocity'].iloc[i]),
                 ha='center', fontsize=9, color='black')

    # Show the plot with value labels
    plt.tight_layout()
    plt.savefig('matrix_velocity_plot.png')  # Saves the plot as a .png file with 300 DPI

    # Reste de la création de la présentation PowerPoint...
    try:
        for slide in prs.slides:
            remplacer_textes(slide,sprintName)

        # Optionnel : Remplacer les images sur les slides spécifiques
        images = ["story_points_pie_chart_image2.png", "image1.png", "matrix_velocity_plot.png"]
        slides_to_modify = [3, 4, 5]  # Indices de diapositive (les indices commencent à 0)

        for i, slide_idx in enumerate(slides_to_modify):
            slide = prs.slides[slide_idx]
            remplacer_image(slide, images[i], image_size=(5, 4))  # Images de 4x3 pouces centrées

        # Enregistrer la présentation modifiée
        prs.save("votre_presentation_modifiee.pptx")

        with tempfile.NamedTemporaryFile(delete=False) as temp_zip:
            with zipfile.ZipFile(temp_zip, 'w') as zipf:
                # Ajouter le fichier Excel à l'archive ZIP
                zipf.writestr(f'vilocity_{sprintName}.xlsx', excel_output.getvalue())
                zipf.write('exelPie.xlsx', arcname=f'{sprintName}_Pie.xlsx')
                zipf.write('votre_presentation_modifiee.pptx', arcname=f'{sprintName}_presentation.pptx')

                # Ajouter le fichier PowerPoint à l'archive ZIP


            temp_zip_path = temp_zip.name

        # Envoyer le fichier PowerPoint
        return send_file(
            temp_zip_path,
            as_attachment=True,
            download_name=f'presentation_and_excel_{file_id}.zip',
            mimetype='application/zip'
        )

    except Exception as e:
        return jsonify({"error": f"Erreur lors de la création de la présentation: {str(e)}"}), 500

@app.route('/data/stats/<int:id>', methods=['GET'])
def get_data_by_id(id):
    connection = get_db_connection()
    if not connection:
        return jsonify({"error": "Failed to connect to the database"}), 500

    try:
        cursor = connection.cursor(dictionary=True)
        query = "SELECT * FROM performanceJira WHERE id = %s"
        cursor.execute(query, (id,))
        result = cursor.fetchone()

        if result:
            # Deserialize JSON fields and encode binary data
            if result['grouped_bar_chart']:
                result['grouped_bar_chart'] = json.loads(result['grouped_bar_chart'])
            if result['pie_chart']:
                result['pie_chart'] = json.loads(result['pie_chart'])
            if result['histogram']:
                result['histogram'] = json.loads(result['histogram'])
            # Encode binary file data to base64
            if result['file_data']:
                result['file_data'] = base64.b64encode(result['file_data']).decode('utf-8')

            return jsonify(result)
        else:
            return jsonify({"error": "Record not found"}), 404
    except Error as e:
        print(f"Error fetching data from MySQL: {e}")
        return jsonify({"error": "Failed to fetch data from the database"}), 500
    finally:
        if connection.is_connected():
            cursor.close()
            connection.close()


from flask import jsonify, request
from mysql.connector import Error


@app.route('/data-perf-jira/<int:id>', methods=['DELETE'])
def delete_data_perf_jira(id):
    connection = get_db_connection()
    if not connection:
        return jsonify({"error": "Failed to connect to the database"}), 500

    try:
        cursor = connection.cursor()

        # Requête pour supprimer l'entrée de la base de données
        query = "DELETE FROM performanceJira WHERE id = %s"
        cursor.execute(query, (id,))
        connection.commit()

        if cursor.rowcount == 0:
            return jsonify({"error": "No DataPerfJira found with this ID"}), 404

        return jsonify({"message": "DataPerfJira deleted successfully"}), 200
    except Error as e:
        print(f"Error deleting data from MySQL: {e}")
        return jsonify({"error": "Failed to delete data from the database"}), 500
    finally:
        if connection.is_connected():
            cursor.close()
            connection.close()


@app.route('/data', methods=['GET'])
def get_data():
    connection = get_db_connection()
    if not connection:
        return jsonify({"error": "Failed to connect to the database"}), 500

    try:
        cursor = connection.cursor(dictionary=True)
        query = "SELECT * FROM performanceJira"
        cursor.execute(query)
        result = cursor.fetchall()

        # Deserialize JSON fields and encode binary data
        for row in result:
            if row['grouped_bar_chart']:
                row['grouped_bar_chart'] = json.loads(row['grouped_bar_chart'])
            if row['pie_chart']:
                row['pie_chart'] = json.loads(row['pie_chart'])
            if row['histogram']:
                row['histogram'] = json.loads(row['histogram'])
            # Encode binary file data to base64
            if row['file_data']:
                row['file_data'] = base64.b64encode(row['file_data']).decode('utf-8')

        return jsonify(result)
    except Error as e:
        print(f"Error fetching data from MySQL: {e}")
        return jsonify({"error": "Failed to fetch data from the database"}), 500
    finally:
        if connection.is_connected():
            cursor.close()
            connection.close()





# Call the function to install the packages


if __name__ == '__main__':

    # List of packages to install
    packages = [
        "Flask", "Flask-Cors", "mysql-connector-python", "paddlepaddle", "numpy", "jira", "pandas", "matplotlib", "pytz"
    ]

    # Install each package using pip
    #for package in packages:
     #   subprocess.check_call([sys.executable, "-m", "pip", "install", package])

    #install_packages()

    # Your main code here
    print("All packages installed successfully!")
    #upgrade_packages()

    # Your main code here
    print("All packages upgraded successfully!")
    create_database_query = "CREATE DATABASE IF NOT EXISTS actiabackImg"
    create_table_query = """
       CREATE TABLE IF NOT EXISTS performanceJira (
        id INT AUTO_INCREMENT PRIMARY KEY,
        start_date DATE,
        end_date DATE,
        sprint_name VARCHAR(255),
        grouped_bar_chart TEXT,
        pie_chart TEXT,
        histogram TEXT,
        file_data LONGBLOB,
        idsprint INT,
        date_genere TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    );

        """

    try:
        # Connect to MySQL without specifying the database
        connection = mysql.connector.connect(
            host=serverdb,
            port='3306',
            user='root',
            password=''
        )

        if connection.is_connected():
            cursor = connection.cursor()

            # Check if the database exists and create it if not
            cursor.execute(create_database_query)
            print("Database 'actiabackImg' checked/created.")

            # Switch to the newly created database
            cursor.execute("USE actiabackImg")

            # Check if the table exists and create it if not
            cursor.execute(create_table_query)
            print("Table 'performanceJira' checked/created.")
    finally:
        if connection.is_connected():
            cursor.close()
            connection.close()
    app.run(host='0.0.0.0', port=5000, debug=True)

    app.run(debug=True)
